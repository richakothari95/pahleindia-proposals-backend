"""
PowerPoint generation service.
Clones slides from PPT_template_V02.pptx and populates them with proposal content.
Uses lxml deepcopy to duplicate the content slide layout (slide 1, 0-indexed).
"""

import os
import copy
from io import BytesIO
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from models.generation import ProposalContent, PPTSlide

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "..", "templates", "PPT_template_V02.pptx")

ORANGE = RGBColor(0xEC, 0x69, 0x1F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _get_shape_by_name(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _set_textbox_text(slide, shape_name: str, text: str):
    shape = _get_shape_by_name(slide, shape_name)
    if shape and shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        # Preserve font size from original if possible
        if tf.paragraphs[0].runs:
            pass  # already set above


def _clone_slide_xml(prs: Presentation, source_index: int) -> etree._Element:
    """Return a deep copy of a slide's XML element."""
    source_slide = prs.slides[source_index]
    return copy.deepcopy(source_slide._element)


def _add_slide_from_xml(prs: Presentation, slide_xml: etree._Element, layout_index: int = 1) -> object:
    """
    Insert a new slide into the presentation using cloned XML.
    Returns the new slide object.
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.oxml.ns import qn
    import random

    slide_layout = prs.slide_layouts[layout_index]
    new_slide = prs.slides.add_slide(slide_layout)

    # Replace the new slide's spTree with the cloned one
    sp_tree = slide_xml.find('.//' + qn('p:spTree'))
    if sp_tree is not None:
        existing_sp_tree = new_slide.shapes._spTree
        existing_parent = existing_sp_tree.getparent()
        idx = list(existing_parent).index(existing_sp_tree)
        existing_parent.remove(existing_sp_tree)
        existing_parent.insert(idx, copy.deepcopy(sp_tree))

    return new_slide


def _remove_slide(prs: Presentation, index: int):
    """Remove a slide by index from the presentation."""
    from pptx.oxml.ns import qn
    xml_slides = prs.slides._sldIdLst
    slides = prs.slides

    # Get the slide element
    slide = slides[index]
    rId = slides._sldIdLst[index].get('r:id')

    # Remove from slide ID list
    xml_slides.remove(xml_slides[index])

    # Remove the relationship
    prs.part.drop_rel(rId)


def _add_content_textbox(slide, body_points: list[str], data_callout: str | None = None):
    """Add a body content text box on the right half of the slide."""
    left = Inches(5.2)
    top = Inches(1.4)
    width = Inches(7.8)
    height = Inches(5.0)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, point in enumerate(body_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"\u2022  {point}"
        p.space_after = Pt(6)
        if p.runs:
            run = p.runs[0]
            run.font.size = Pt(16)
            run.font.color.rgb = WHITE
            run.font.name = "Lato"

    if data_callout:
        p = tf.add_paragraph()
        p.space_before = Pt(12)
        p.text = data_callout
        if p.runs:
            run = p.runs[0]
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = ORANGE
            run.font.name = "Lato"


class PptxService:
    def generate(self, content: ProposalContent) -> bytes:
        prs = Presentation(TEMPLATE_PATH)

        # Save content slide XML before removing template slides
        # Template: [0]=Title, [1]=Content, [2]=Content, [3]=ThankYou
        content_slide_xml = _clone_slide_xml(prs, 1)

        # Update Slide 0 (Title slide) with proposal title
        title_slide = prs.slides[0]
        for shape in title_slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Find the title placeholder text box (largest text area)
                        if len(run.text) > 5 and "Pahle" not in run.text and "www" not in run.text:
                            run.text = content.title
                            break

        # Remove template content slides (indices 2 and 1, in reverse order to preserve indexing)
        try:
            _remove_slide(prs, 2)
        except Exception:
            pass
        try:
            _remove_slide(prs, 1)
        except Exception:
            pass

        # Now prs.slides = [Title slide (0), Thank You slide (1)]
        # We insert content slides before the Thank You slide

        slides_to_add = content.ppt_slides
        # Filter out title and thank_you from the generation list
        content_slides = [s for s in slides_to_add if s.slide_type not in ("title", "thank_you")]

        for i, slide_def in enumerate(content_slides):
            # Insert at position len(prs.slides) - 1 (before Thank You)
            # We use add_slide and then move it
            new_slide = _add_slide_from_xml(prs, content_slide_xml, layout_index=1)

            # Set title (TextBox 3)
            _set_textbox_text(new_slide, "TextBox 3", slide_def.title)

            # Set subtitle (TextBox 7)
            if slide_def.subtitle:
                _set_textbox_text(new_slide, "TextBox 7", slide_def.subtitle)
            else:
                _set_textbox_text(new_slide, "TextBox 7", "")

            # Set slide number (TextBox 12)
            _set_textbox_text(new_slide, "TextBox 12", str(i + 2).zfill(2))

            # Add body content
            if slide_def.body_points:
                _add_content_textbox(new_slide, slide_def.body_points, slide_def.data_callout)

        # Move the Thank You slide to the end
        # After add_slide calls, Thank You is at index 1; content slides were appended after it
        # Re-order: move slide at index 1 (Thank You) to the end
        try:
            from pptx.oxml.ns import qn
            sldIdLst = prs.slides._sldIdLst
            # The Thank You slide ID element is at index 1
            thank_you_elem = sldIdLst[1]
            sldIdLst.remove(thank_you_elem)
            sldIdLst.append(thank_you_elem)
        except Exception:
            pass

        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()
